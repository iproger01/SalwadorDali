import example as example
from pptx import Presentation
from datetime import date
from pptx.util import Cm, Pt, Inches
from dateutil.parser import parse
import asyncio

root = Presentation()

slide_layout = root.slide_layouts[0]

""" Ref for slide types: 
0 ->  title and subtitle
1 ->  title and content
2 ->  section header
3 ->  two content
4 ->  Comparison
5 ->  Title only 
6 ->  Blank
7 ->  Content with caption
8 ->  Pic with caption
"""
def generate_abonement(buyer):
    slide = root.slides.add_slide(slide_layout)
    left = top = 0
    pic = slide.shapes.add_picture('abonement.jpg', left, top, width=root.slide_width, height = root.slide_height)

    # slide0 = root.slides[0]
    title_name=slide.shapes.title
    title_name.top = Cm(7.8)
    title_name.left = Cm(2)
    title_name.text = buyer

    dt_now = str(date.today())
    dt_parse=parse(dt_now).strftime('%d.%m.%Y')

    left = Cm(20)
    top = Cm(17)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = dt_parse
    # date_of_buy.top = Cm(18)
    # date_of_buy.left = Cm(19)

    # date_of_buy = slide.placeholders[1]

    # date_of_buy.text = dt_parse
    # shape.text_frame.paragraphs[0].font.size = Pt(15)
    # date_of_buy.font.size = Pt(14)

    slide.shapes._spTree.insert(2, pic._element)
    root.save(f"{buyer}.pptx")

if __name__ == "__main__":
    generate_abonement('Natali Ustimenko')